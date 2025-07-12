from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from pptx import Presentation
from pptx.util import Inches, Pt
import textwrap

def get_bible_text(bible_json:dict, bible_head_dict:dict, head_str:str, f_num:str, s_start_num:str, s_end_num:str = '0') -> list:
    # 성경 리스트로 불러오기
    the_text_li = []
    the_text_li.append(s_start_num+'. '+bible_json[head_str][f_num][s_start_num])
    if s_end_num == '0':
        return the_text_li
    else:
        for idx in range(int(s_start_num)+1, int(s_end_num)+1):
            the_text_li.append(str(idx)+'. '+bible_json[head_str][f_num][str(idx)])

    # ppt 파일 이름 생성
    ppt_file_nm = bible_head_dict[head_str]+f_num+'_'+s_start_num
    if s_end_num!= '0':
        ppt_file_nm = ppt_file_nm +'-'+s_end_num
    
    return the_text_li, ppt_file_nm

# the_text_li, ppt_file_nm = get_bible_text(bible_json, '창', '1', '1', '4')

def add_text_to_textbox(text, text_box_size, font_size):
    """텍스트 상자에 텍스트를 추가하고 줄바꿈을 처리합니다.

    Args:
        textbox: 텍스트 상자 객체.
        text: 추가할 텍스트 문자열.
        width_ratio: 텍스트 상자 너비 대비 문자 수 비율 (조절 필요).
    """
    wrapped_text = textwrap.wrap(text, width=int(Inches(text_box_size) / (Pt(font_size)*0.9)))
    return "\n    ".join(wrapped_text), len(wrapped_text)

def get_bible_ppt(image_files_path, the_text_li, ppt_file_nm, ppt_title):
    # 이미지와 텍스트리스트, 파일명을 받아 ppt로 구성해서 저장하는 함수

    # 프레젠테이션 객체 생성
    prs = Presentation()
    text_box_width = 13
    text_box_height = 6
    font_size = 36
    lined_text = ''
    line_count = 0
    max_line_count = int(Inches(text_box_height) / (Pt(font_size)*1.7))
    print(max_line_count)

    if ppt_title == '':
            ppt_title = ':'.join(ppt_file_nm.split('_'))
            sub_title = ''
    else:
        sub_title = ':'.join(ppt_file_nm.split('_'))

    # 각 이미지에 대해 슬라이드 추가
    for idx in range(len(the_text_li)):
        the_text = the_text_li[idx]

        # 줄바꿈 처리
        now_lined_text, now_line_count = add_text_to_textbox(the_text, text_box_size=text_box_width, font_size=font_size)
        lined_text = '\n'.join([lined_text, now_lined_text]).strip()
        line_count+=now_line_count

        # 최대 절수 넣는 판단
        if idx+1 < len(the_text_li):
            next_lined_text, next_line_count = add_text_to_textbox(the_text_li[idx+1], text_box_size=text_box_width, font_size=font_size)    
            if line_count+next_line_count < max_line_count:
                continue
            else:
                pass
        else:
            pass

        # 슬라이드 추가
        prs.slide_width = Inches(16) 
        prs.slide_height = Inches(9)
        slide_layout = prs.slide_layouts[1]  # 빈 슬라이드 레이아웃 선택
        slide = prs.slides.add_slide(slide_layout)

        # 이미지 추가
        left = Inches(0)
        top = Inches(0)
        height = Inches(9)
        
        slide.shapes.add_picture(image_files_path, left, top, width=Inches(16), height=Inches(9))
        # slide.shapes.add_picture(image_files, left, top, height=height)

        # 제목
        # shape = slide.shapes.title
        # shape.text = ':'.join(ppt_file_nm.split('_'))


        title_box = slide.shapes.add_textbox(left=Inches(1), top=Inches(0.2), width=Inches(14), height=Inches(4))
        title_box_tf = title_box.text_frame
        title_box_tf_p = title_box_tf.add_paragraph()
        title_box_tf_p.text = ppt_title
        title_box_tf_p.font.size = Pt(43)
        title_box_tf_p.font.bold = True
        try:
           title_box_tf_p.font.name = '맑은 고딕'
        except:
           pass

        # 서브 제목
        # shape = slide.shapes.title
        # shape.text = ':'.join(ppt_file_nm.split('_'))


        title_box = slide.shapes.add_textbox(left=Inches(1), top=Inches(1.0), width=Inches(14), height=Inches(4))
        title_box_tf = title_box.text_frame
        title_box_tf_p = title_box_tf.add_paragraph()
        title_box_tf_p.text = sub_title
        title_box_tf_p.font.size = Pt(30)
        # title_box_tf_p.font.bold = True
        try:
           title_box_tf_p.font.name = '맑은 고딕'
        except:
           pass

        # 내용
        text_box = slide.shapes.add_textbox(left=Inches(1), top=Inches(1.75), width=Inches(14), height=Inches(6))
        text_box_tf = text_box.text_frame
        text_box_tf_p = text_box_tf.add_paragraph()
        text_box_tf_p.text = lined_text 
        text_box_tf_p.font.size = Pt(40)
        try:
            text_box_tf_p.font.name = '넥슨 풋볼고딕 B'
            # text_box_tf_p.font.size = Pt(40)
        except:
            pass

        lined_text = ''
        line_count = 0
        
        #tf = body_shape.text_frame
        #tf.text = the_text

    # 파일 저장
    
    prs.save('./ppts/'+ppt_file_nm+'.pptx')



import json

with open('./data/index_bible2.json', 'r', encoding='utf8') as b_json:
    bible_json = json.load(b_json)

with open('./data/bible_head_dict.json', 'r', encoding='utf8') as b_json:
    bible_head_dict = json.load(b_json)
    
  # 여기에 이미지 파일 경로를 추가하세요
image_files_path = './data/ppt배경.png'


from tkinter import *
import tkinter.ttk
import re

search_content = ''
target_book_nm = '창'
target_chapter_num = '1'
target_start_line = '1'
chapter_li = []
title = ''

root = Tk()
root.title('BIBLE2PPT')
root.geometry("500x500")

# etc_frame = Frame(root, relief='solid', bd=2, background='gray', height = 100)
# etc_frame.pack(fill='both', expand=True)

select_bible_frame = Frame(root, relief='solid', bd=2, background='gray', height = 100)
select_bible_frame.pack(fill='both', expand=True)
# select_bible_lable = Label(select_bible_frame, text='성경 선택', foreground='white', background='gray')
# select_bible_lable.grid(column = 1, row = 1)


# 드롭박스로 성경 선택
def update_chapters(eventObject): # eventObejct 자리에는 아무 값이나 들어가도 괜찮습니다.
    # book_combobox에 의해 chapter_combobox가 반응하여 바뀌도록 해주는 역할
    global target_book_nm
    if book_combobox.get() != target_book_nm:
        target_book_nm = book_combobox.get()
        chapter_combobox.config(value=list(bible_json[target_book_nm].keys()))
        chapter_combobox.set('장')
        line1_combobox.config(value=list(bible_json[target_book_nm][target_chapter_num].keys()))
        line1_combobox.set('절_시작')
        line2_combobox.config(value=list(bible_json[target_book_nm][target_chapter_num].keys()))
        line2_combobox.set('절_끝')

def update_line(eventObject): # eventObejct 자리에는 아무 값이나 들어가도 괜찮습니다.
    # chapter_combobox에 의해 line_combobox가 반응하여 바뀌도록 해주는 역할
    global target_book_nm, target_chapter_num
    if chapter_combobox.get() != target_chapter_num:
        target_chapter_num = chapter_combobox.get()
        line1_combobox.config(value=list(bible_json[target_book_nm][target_chapter_num].keys()))
        line1_combobox.set('절_시작')
        line2_combobox.config(value=list(bible_json[target_book_nm][target_chapter_num].keys()))
        line2_combobox.set('절_끝')

def update_line2(eventObject): # eventObejct 자리에는 아무 값이나 들어가도 괜찮습니다.
    # chapter_combobox에 의해 line_combobox가 반응하여 바뀌도록 해주는 역할
    global target_book_nm, target_chapter_num, target_start_line 
    if line1_combobox.get() != target_start_line:
        target_start_line = line1_combobox.get()
        line2_combobox.config(value=list(bible_json[target_book_nm][target_chapter_num].keys())[int(target_start_line):])
        line2_combobox.set('절_끝')

def get_dropbox_search():
    global search_content 
    book = book_combobox.get()
    chapter = chapter_combobox.get()
    line_start = line1_combobox.get()
    line_end = line2_combobox.get()
    search_content = book+str(chapter)+':'+str(line_start)+ '-' +str(line_end)

drop_down_lable = Label(select_bible_frame, text='성경선택', foreground='white', background='gray')
drop_down_lable.grid(column = 1, row = 2)
book_lable = Label(select_bible_frame, text='성경', foreground='white', background='gray')
book_lable.grid(column = 2, row = 2)
book_combobox = tkinter.ttk.Combobox(select_bible_frame, height=7, values=list(bible_json.keys()))
book_combobox.grid(column = 3, row = 2)
book_combobox.set('성경')
book_combobox.bind('<<ComboboxSelected>>', update_chapters) # book_combobox에 의해 chapter_combobox가 반응하여 바뀌도록 해주는 역할

chapter_lable = Label(select_bible_frame, text='장', foreground='white', background='gray')
chapter_lable.grid(column = 2, row = 3)
chapter_combobox = tkinter.ttk.Combobox(select_bible_frame, height=7, values=chapter_li)
chapter_combobox.grid(column = 3, row = 3)
chapter_combobox.set('장')
chapter_combobox.bind('<<ComboboxSelected>>', update_line) # book_combobox에 의해 chapter_combobox가 반응하여 바뀌도록 해주는 역할

line1_lable = Label(select_bible_frame, text='절_시작', foreground='white', background='gray')
line1_lable.grid(column = 2, row = 4)
line1_combobox = tkinter.ttk.Combobox(select_bible_frame, height=7, values=chapter_li)
line1_combobox.grid(column = 3, row = 4)
line1_combobox.set('절_시작')
line1_combobox.bind('<<ComboboxSelected>>', update_line2)

line2_lable = Label(select_bible_frame, text='절_끝', foreground='white', background='gray')
line2_lable.grid(column = 2, row = 5)
line2_combobox = tkinter.ttk.Combobox(select_bible_frame, height=7, values=chapter_li)
line2_combobox.grid(column = 3, row = 5)
line2_combobox.set('절_끝')

search_button = Button(select_bible_frame, text = '검색', command=get_dropbox_search)
search_button.grid(column=3, row=6)


# 직접 입력
direct_search_frame = Frame(root, relief='solid', bd=2, background='gray', height = 100)
direct_search_frame.pack(fill='both', expand=True)
# direct_search_lable = Label(direct_search_frame, text='직접 검색', foreground='white', background='gray')
# direct_search_lable.grid(column = 1, row = 1)

def get_direct_search():
    global search_content 
    search_content = direct_search_entry.get()

direct_search_lable = Label(direct_search_frame, text='직접검색(예, 창14:1-5)', foreground='white', background='gray')
direct_search_lable.grid(column = 1, row = 6)

direct_search_entry = Entry(direct_search_frame, foreground='black', background='white')
direct_search_entry.insert(0, '')
direct_search_entry.grid(column=2, row=6)

search_button = Button(direct_search_frame, text = '검색', command=get_direct_search)
search_button.grid(column=3, row=6)


# ppt 제목 입력
def get_title():
    global title 
    title = insert_title_entry.get()

insert_title_frame = Frame(root, relief='solid', bd=2, background='gray', height = 100)
insert_title_frame.pack(fill='both', expand=True)
insert_title_lable = Label(insert_title_frame, text='제목 입력', foreground='white', background='gray')
insert_title_lable.grid(column = 1, row = 1)

insert_title_entry = Entry(insert_title_frame, foreground='black', background='white')
insert_title_entry.insert(0, '')
insert_title_entry.grid(column=2, row=1)

insert_title_button = Button(insert_title_frame, text = '입력', command=get_title)
insert_title_button.grid(column=3, row=1)


# ppt 생성
def get_ppt():
    global search_content, image_files_path, bible_json, bible_head_dict, title

    first_step = re.findall('[0-9]*:[0-9]*-[0-9]*', search_content)[0]
    book = search_content.split(first_step)[0]
    chapter, lines = first_step.split(':')
    line_start, line_end = lines.split('-')

    the_text_li, ppt_file_nm = get_bible_text(bible_json, bible_head_dict, book, chapter, line_start, line_end)
    get_bible_ppt(image_files_path, the_text_li, ppt_file_nm, title)

get_ppt_button_frame = Frame(root, relief='solid', bd=2, background='gray', height = 100)
get_ppt_button_frame.pack(fill='both', expand=True)

get_ppt_button = Button(get_ppt_button_frame, text = '생성', command=get_ppt)
get_ppt_button.grid(column=3, row=1)



root.mainloop()