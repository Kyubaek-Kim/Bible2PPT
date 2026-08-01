"""Slide engine: turns aligned verse bundles into a ``.pptx``.

Responsibilities:

* aspect ratio → slide size and text-box geometry (16:9 / 4:3 / A4);
* fonts → title / section-info / body, each with its own family / size / bold
  (body weight can be forced via the UI checkbox), East-Asian names set
  correctly; title and over-long lines auto-shrink to fit;
* optional user layout overrides → title / section / body boxes stored as
  fractional [x, y, w, h] rectangles so a saved layout is aspect-independent;
* full-slide background image;
* verse formatting → verse-number prefix, one verse per line, ``<N장>`` marker at
  chapter changes, cross-translation interleave;
* pagination → a verse *bundle* (all translations of one canonical verse) never
  splits across slides; lines-per-slide are recomputed from font/size/ratio. The
  body is packed to use the available area well, but never overflows the slide
  nor the header band: if an indivisible bundle would overflow, the body font is
  shrunk (down to a safe minimum) until it fits, and only if even that is
  impossible is a structured :class:`PaginationError` raised — the engine never
  silently clips or overflows.
"""
from __future__ import annotations

import re
import unicodedata
from dataclasses import dataclass, field, replace
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from . import fonts
from .alignment import VerseBundle

# Aspect ratio -> (width_in, height_in).
ASPECT_RATIOS: dict[str, tuple[float, float]] = {
    "16:9": (13.333, 7.5),
    "4:3": (10.0, 7.5),
    "A4": (11.69, 8.27),  # landscape A4
}
DEFAULT_ASPECT = "16:9"

TITLE_FONT_SIZE = 40
SECTION_FONT_SIZE = 26
MARGIN_IN = 0.6
# body text sits in a roomier box than the title/section so verses have more
# breathing space on the left, right and bottom of the slide.
BODY_SIDE_MARGIN_IN = 0.95
BODY_BOTTOM_MARGIN_IN = 0.7
IN_TO_PT = 72.0
LINE_SPACING = 1.3
# A font's intrinsic single-line height is larger than its point size; the
# "multiple" line-spacing above is applied on top of it. Fold both into the
# vertical-fit estimate so bundles don't overflow the bottom of a slide. This is
# a *safe upper bound* over the win-ascent/descent of the bundled faces
# (NanumGothic ≈1.08, NanumSquare ≈1.11), so the fit never overflows while still
# packing verses more tightly than the previous 1.2 guess.
FONT_LINE_HEIGHT = 1.15
# where the body starts (inches from the top) when the header (title/section) is
# present vs. hidden — hiding both reclaims that space for more verses.
BODY_TOP_WITH_HEADER_IN = 2.0
BODY_TOP_NO_HEADER_IN = 0.5
# The body hanging indent is sized per passage from the widest verse number
# (see ``body_hang_pt``): the number hangs left of a tab stop where every text
# line begins, giving a straight left edge regardless of number width.
# how far the body font size may be auto-reduced to fit verses tidily
MAX_BODY_SHRINK_PT = 2
MIN_BODY_FONT_SIZE = 12
# title auto-shrink bounds (keep the title on one line within the slide width)
MIN_TITLE_FONT_SIZE = 18


def _display_units(text: str) -> int:
    """Approximate rendered width in half-em units (CJK≈2, latin≈1)."""
    total = 0
    for ch in text:
        total += 2 if unicodedata.east_asian_width(ch) in ("W", "F") else 1
    return total


@dataclass
class SlideStyle:
    aspect: str = DEFAULT_ASPECT
    font_name: str = "나눔스퀘어 Bold"  # stored dropdown label
    body_font_size: int = 32
    # ``None`` -> follow the font face's intrinsic weight; a bool is an explicit
    # user choice from the body "bold" checkbox.
    body_bold_opt: bool | None = None
    # Per-element title / section (reference) styling. The *font face* is always
    # the single global ``font_name`` (chosen in 화면 설정); only size / bold /
    # visibility differ per element, so a title can never get stuck on a stale
    # per-element face.
    title_font_size: int = TITLE_FONT_SIZE
    title_bold: bool = True
    title_enabled: bool = True
    section_font_size: int = SECTION_FONT_SIZE
    section_bold: bool = True
    section_enabled: bool = True
    # Text colours as ``#rrggbb`` (empty -> engine default black).
    title_color: str = ""
    section_color: str = ""
    body_color: str = ""
    # Fractional [x, y, w, h] box overrides keyed by "title"/"section"/"body".
    layout_boxes: dict[str, list[float]] = field(default_factory=dict)

    @property
    def _choice(self) -> fonts.FontChoice:
        return fonts.resolve(self.font_name)

    @property
    def typeface(self) -> str:
        """Actual font family to set on runs (resolved from the label).

        Every element (title / section / body) shares this single family.
        """
        return self._choice.typeface

    # title / section share the global body typeface (no per-element face)
    title_typeface = typeface
    section_typeface = typeface

    @property
    def body_bold(self) -> bool:
        """Body weight: explicit user choice, else the face's intrinsic weight."""
        if self.body_bold_opt is not None:
            return self.body_bold_opt
        return self._choice.bold

    def run_bold(self, user_bold: bool) -> bool:
        """Effective bold bit for a run in the global face (see ``fonts.run_bold``)."""
        return fonts.run_bold(self._choice, user_bold)

    @property
    def size_in(self) -> tuple[float, float]:
        return ASPECT_RATIOS.get(self.aspect, ASPECT_RATIOS[DEFAULT_ASPECT])

    def _override(self, key: str) -> tuple[float, float, float, float] | None:
        frac = self.layout_boxes.get(key)
        if not frac or len(frac) != 4:
            return None
        w, h = self.size_in
        x, y, bw, bh = frac
        return x * w, y * h, bw * w, bh * h

    def _default_title_box(self) -> tuple[float, float, float, float]:
        w, _ = self.size_in
        return MARGIN_IN, 0.25, w - 2 * MARGIN_IN, 1.1

    def _default_section_box(self) -> tuple[float, float, float, float]:
        w, _ = self.size_in
        return MARGIN_IN, 1.15, w - 2 * MARGIN_IN, 0.8

    def _default_body_box(self) -> tuple[float, float, float, float]:
        w, h = self.size_in
        # reclaim the header band for the body when both title and section are
        # hidden, so verses can use (almost) the whole slide.
        top = BODY_TOP_WITH_HEADER_IN
        if not self.title_enabled and not self.section_enabled:
            top = BODY_TOP_NO_HEADER_IN
        return BODY_SIDE_MARGIN_IN, top, w - 2 * BODY_SIDE_MARGIN_IN, h - top - BODY_BOTTOM_MARGIN_IN

    @property
    def body_box(self) -> tuple[float, float, float, float]:
        """left, top, width, height (inches) of the body text area."""
        return self._override("body") or self._default_body_box()

    @property
    def title_box(self) -> tuple[float, float, float, float]:
        return self._override("title") or self._default_title_box()

    @property
    def section_box(self) -> tuple[float, float, float, float]:
        return self._override("section") or self._default_section_box()

    def default_layout_fractions(self) -> dict[str, list[float]]:
        """Engine-default boxes as fractional [x, y, w, h] of the slide.

        Used by the layout-customisation UI to seed the draggable rectangles
        (and by "reset")."""
        w, h = self.size_in
        out: dict[str, list[float]] = {}
        for key, box in (
            ("title", self._default_title_box()),
            ("section", self._default_section_box()),
            ("body", self._default_body_box()),
        ):
            x, y, bw, bh = box
            out[key] = [x / w, y / h, bw / w, bh / h]
        return out

    @property
    def max_units_per_line(self) -> int:
        _, _, width_in, _ = self.body_box
        width_pt = width_in * IN_TO_PT
        # each half-em unit ~ font_size * 0.5 pt wide
        return max(1, int(width_pt / (self.body_font_size * 0.5)))

    @property
    def line_pitch_pt(self) -> float:
        """Vertical distance between consecutive body lines, in points."""
        return self.body_font_size * LINE_SPACING * FONT_LINE_HEIGHT

    @property
    def max_body_lines(self) -> int:
        _, _, _, height_in = self.body_box
        height_pt = height_in * IN_TO_PT
        return max(1, int(height_pt / self.line_pitch_pt))

    def hang_units(self, hang_pt: float) -> int:
        """Hanging-indent width expressed in half-em units for the body size."""
        return max(0, round(hang_pt / (self.body_font_size * 0.5)))


def wrap_line(text: str, max_units: int, cont_units: int | None = None) -> list[str]:
    """Word-wrap ``text`` to ``max_units`` half-em units, hard-splitting long tokens.

    When ``cont_units`` is given, the first line is capped at ``max_units`` and
    every subsequent (wrapped) line at ``cont_units`` — this models the hanging
    indent, where wrapped lines are narrower than the first.
    """
    if not text:
        return [""]

    def cap(n_done: int) -> int:
        return max_units if (n_done == 0 or cont_units is None) else cont_units

    lines: list[str] = []
    current = ""
    current_units = 0
    for token in re.split(r"(\s+)", text):
        if token == "":
            continue
        tu = _display_units(token)
        if token.isspace():
            if current:
                current += token
                current_units += tu
            continue
        limit = cap(len(lines))
        if current_units + tu <= limit or not current:
            if tu > limit and not current:
                # token itself longer than a line -> hard split
                for chunk in _hard_split(token, limit):
                    lines.append(chunk)
                current, current_units = "", 0
                continue
            current += token
            current_units += tu
        else:
            lines.append(current.rstrip())
            limit = cap(len(lines))
            if tu > limit:
                for chunk in _hard_split(token, limit):
                    lines.append(chunk)
                current, current_units = "", 0
            else:
                current, current_units = token, tu
    if current.strip():
        lines.append(current.rstrip())
    return lines or [""]


def _hard_split(token: str, max_units: int) -> list[str]:
    out: list[str] = []
    buf = ""
    units = 0
    for ch in token:
        cu = _display_units(ch)
        if units + cu > max_units and buf:
            out.append(buf)
            buf, units = ch, cu
        else:
            buf += ch
            units += cu
    if buf:
        out.append(buf)
    return out


def meaningful_title(title: str) -> bool:
    """True if the title contains any letter/number/CJK (not blank/punct-only)."""
    return bool(re.search(r"[0-9A-Za-z\uAC00-\uD7A3\u3040-\u30FF\u4E00-\u9FFF]", title or ""))


# --------------------------------------------------------------------------- #
# Pagination
# --------------------------------------------------------------------------- #
@dataclass
class RenderLine:
    text: str
    kind: str  # "chapter" | "verse"


@dataclass
class SlidePage:
    lines: list[RenderLine] = field(default_factory=list)


def _bundle_block(bundle: VerseBundle, with_marker: bool) -> list[RenderLine]:
    """Flatten one bundle to render lines (optional chapter marker + interleaved cells)."""
    lines: list[RenderLine] = []
    if with_marker:
        # leading tab aligns the marker with the verse-text column (marL)
        lines.append(RenderLine(text=f"\t<{bundle.coord.chapter}장>", kind="chapter"))
    for _code, cell in bundle.cells:
        if not cell.visible:
            continue
        # tab after the number: the number sits in the outdented column while the
        # text starts at the tab stop, keeping the body's left edge straight.
        lines.append(RenderLine(text=f"{cell.label}.\t{cell.text}", kind="verse"))
    return lines


def _blocks(bundles: list[VerseBundle]) -> list[list[RenderLine]]:
    """One indivisible render block per bundle; chapter markers decided once.

    A ``<N장>`` marker is attached to a bundle iff the passage spans multiple
    chapters and this bundle opens a new chapter (reading order). Markers are
    never duplicated across page breaks — they mark true change points only.
    """
    multi_chapter = len({b.coord.chapter for b in bundles}) > 1
    blocks: list[list[RenderLine]] = []
    prev_chapter: int | None = None
    for bundle in bundles:
        with_marker = multi_chapter and bundle.coord.chapter != prev_chapter
        blocks.append(_bundle_block(bundle, with_marker))
        prev_chapter = bundle.coord.chapter
    return blocks


def _measure_text_pt(text: str, size: int, style: SlideStyle) -> float:
    """Approximate rendered width of ``text`` in points at ``size``.

    Uses the bundled TTF via Pillow when available (exact); otherwise falls back
    to a per-character estimate (digits are wider than punctuation)."""
    choice = fonts.resolve(style.font_name)
    if choice.bundled and choice.bundled.exists():
        try:
            from PIL import ImageFont

            font = ImageFont.truetype(str(choice.bundled), size)
            return float(font.getlength(text))  # px at 72dpi ≈ points
        except Exception:  # noqa: BLE001 - fall back to the estimate
            pass
    width = 0.0
    for ch in text:
        if ch.isdigit():
            width += 0.60 * size
        elif ch in ".:-–—":
            width += 0.32 * size
        elif unicodedata.east_asian_width(ch) in ("W", "F"):
            width += 1.0 * size
        else:
            width += 0.5 * size
    return width


def body_hang_pt(bundles: list[VerseBundle], style: SlideStyle) -> float:
    """Left outdent (points) sized to the widest verse number in the passage.

    Every body line's text begins here (via a tab stop), so a passage whose
    numbers are all 1–2 digits gets a tight indent while a Psalm 119 passage
    (up to ``176``) gets a wider one — in both cases the text column is uniform.
    """
    labels = [cell.label for b in bundles for _c, cell in b.cells if cell.visible]
    size = style.body_font_size
    widest = max((_measure_text_pt(f"{lb}.", size, style) for lb in labels), default=size)
    return widest + 0.45 * size  # gap between the number and the text


def _text_units(style: SlideStyle, hang_pt: float) -> int:
    """Usable text width (half-em units) after the hanging indent."""
    return max(1, style.max_units_per_line - style.hang_units(hang_pt))


def _block_line_count(block: list[RenderLine], style: SlideStyle, hang_pt: float) -> int:
    text_units = _text_units(style, hang_pt)
    return sum(len(wrap_line(ln.text, text_units)) for ln in block)


def fit_body_style(bundles: list[VerseBundle], style: SlideStyle) -> SlideStyle:
    """Auto-reduce the body font (up to ``MAX_BODY_SHRINK_PT``) so verses fit tidily.

    Picks the largest size within the shrink budget for which no single bundle
    overflows a slide, which both avoids downward overflow and reduces awkward
    single-verse-per-slide splits. If nothing fits, uses the smallest tried.
    """
    base = style.body_font_size
    floor = max(MIN_BODY_FONT_SIZE, base - MAX_BODY_SHRINK_PT)
    blocks = _blocks(bundles)
    smallest = style
    for size in range(base, floor - 1, -1):
        trial = replace(style, body_font_size=size)
        smallest = trial
        hang = body_hang_pt(bundles, trial)
        if all(_block_line_count(b, trial, hang) <= trial.max_body_lines for b in blocks):
            return trial
    return smallest


class PaginationError(Exception):
    """Raised when an indivisible verse bundle cannot fit a slide even at the
    minimum body font size (so the engine never silently overflows)."""


def paginate(bundles: list[VerseBundle], style: SlideStyle, hang_pt: float | None = None) -> list[SlidePage]:
    """Greedy pagination that keeps each bundle intact on a single slide."""
    max_lines = style.max_body_lines
    if hang_pt is None:
        hang_pt = body_hang_pt(bundles, style)

    pages: list[SlidePage] = []
    current = SlidePage()
    current_count = 0

    for block in _blocks(bundles):
        wrapped = _block_line_count(block, style, hang_pt)
        if current.lines and current_count + wrapped > max_lines:
            pages.append(current)
            current = SlidePage()
            current_count = 0
        current.lines.extend(block)
        current_count += wrapped
        # a bundle that alone exceeds a slide is isolated on its own page; the
        # overflow is then eliminated by shrinking in ``fit_pages`` (or reported).
        if current_count > max_lines and current.lines == block:
            pages.append(current)
            current = SlidePage()
            current_count = 0

    if current.lines:
        pages.append(current)
    return pages


def _page_fits(page: SlidePage, style: SlideStyle, hang_pt: float) -> bool:
    return _block_line_count(page.lines, style, hang_pt) <= style.max_body_lines


def fit_pages(
    bundles: list[VerseBundle], style: SlideStyle,
) -> tuple[SlideStyle, list[SlidePage], float]:
    """Paginate ``bundles`` guaranteeing that *no* page overflows the body area.

    Starts from ``style`` (already tidied by :func:`fit_body_style`) and, only
    if some indivisible bundle still overflows, shrinks the body font further —
    down to :data:`MIN_BODY_FONT_SIZE` — re-paginating each step. Returns the
    style/pages/hanging-indent actually used. Raises :class:`PaginationError`
    if even the minimum size cannot contain a bundle.
    """
    for size in range(style.body_font_size, MIN_BODY_FONT_SIZE - 1, -1):
        trial = replace(style, body_font_size=size)
        hang = body_hang_pt(bundles, trial)
        pages = paginate(bundles, trial, hang)
        if all(_page_fits(pg, trial, hang) for pg in pages):
            return trial, pages, hang
    raise PaginationError(
        "본문 한 절이 너무 길어 최소 글자 크기로도 한 장표에 담을 수 없습니다."
    )


# --------------------------------------------------------------------------- #
# Rendering
# --------------------------------------------------------------------------- #
@dataclass
class PassageContent:
    title: str
    section_info: str
    bundles: list[VerseBundle]


def _hex_rgb(color: str) -> RGBColor | None:
    """Parse ``#rrggbb`` / ``rrggbb`` to an ``RGBColor`` (None if blank/invalid)."""
    if not color:
        return None
    s = color.lstrip("#").strip()
    if len(s) != 6:
        return None
    try:
        return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except ValueError:
        return None


def _set_font(run, name: str, size: int, *, bold: bool = False, color: str = "") -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = name
    rgb = _hex_rgb(color)
    if rgb is not None:
        run.font.color.rgb = rgb
    # ensure East-Asian text also uses the chosen family
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def _set_hanging_indent(p, hang_pt: float) -> None:
    """Outdent the first line by ``hang_pt`` so the verse number sticks out left
    and wrapped continuation lines align under the verse text.

    A left tab stop is placed at the text-start position (``hang``). Body lines
    put a tab right after the verse number, so the number sits alone in the
    outdented column while *every* text line — first and wrapped alike — begins
    at ``hang``. This makes the body's left edge a straight vertical line instead
    of zig-zagging between the number column and the wrapped-line column."""
    hang = int(Pt(hang_pt))
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(hang))
    pPr.set("indent", str(-hang))
    for existing in pPr.findall(qn("a:tabLst")):
        pPr.remove(existing)
    tab_lst = pPr.makeelement(qn("a:tabLst"), {})
    tab_lst.append(pPr.makeelement(qn("a:tab"), {"pos": str(hang), "algn": "l"}))
    pPr.append(tab_lst)


def _add_textbox(
    slide, box, text_lines, name, size, *,
    bold=False, color="", align=PP_ALIGN.LEFT, line_spacing=None, hanging_pt=None,
):
    left, top, width, height = box
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in text_lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        if hanging_pt:
            _set_hanging_indent(p, hanging_pt)
        run = p.add_run()
        run.text = line
        _set_font(run, name, size, bold=bold, color=color)
    return tb


def _fit_single_line_size(text: str, box_width_in: float, base_size: int, min_size: int) -> int:
    """Largest size (≤ ``base_size``) at which ``text`` fits the box width on one line."""
    width_pt = box_width_in * IN_TO_PT
    units = _display_units(text)
    if units == 0:
        return base_size
    for size in range(base_size, min_size - 1, -1):
        if units * (size * 0.5) <= width_pt:
            return size
    return min_size


def _render_page(prs, page: SlidePage, passage: PassageContent, style: SlideStyle,
                 background: Path | None, hang_pt: float):
    w, h = style.size_in
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank)

    if background is not None and Path(background).exists():
        slide.shapes.add_picture(
            str(background), Inches(0), Inches(0), width=Inches(w), height=Inches(h)
        )

    title_w = style.title_box[2]
    has_title = meaningful_title(passage.title) and style.title_enabled
    if has_title:
        # shrink an over-long title so it never runs past the slide edge
        title_size = _fit_single_line_size(passage.title, title_w, style.title_font_size, MIN_TITLE_FONT_SIZE)
        _add_textbox(slide, style.title_box, [passage.title], style.title_typeface,
                     title_size, bold=style.run_bold(style.title_bold), color=style.title_color)
        if style.section_enabled:
            _add_textbox(slide, style.section_box, [passage.section_info], style.section_typeface,
                         style.section_font_size, bold=style.run_bold(style.section_bold),
                         color=style.section_color)
    elif style.section_enabled:
        # blank/disabled title -> put section info in the title position (task 13)
        sec_size = _fit_single_line_size(passage.section_info, title_w, style.title_font_size, MIN_TITLE_FONT_SIZE)
        _add_textbox(slide, style.title_box, [passage.section_info], style.title_typeface,
                     sec_size, bold=style.run_bold(style.title_bold), color=style.section_color)

    body_lines = [ln.text for ln in page.lines]
    _add_textbox(
        slide, style.body_box, body_lines, style.typeface, style.body_font_size,
        bold=style.run_bold(style.body_bold),
        color=style.body_color,
        line_spacing=LINE_SPACING,
        hanging_pt=hang_pt,
    )


def render(
    passages: list[PassageContent],
    style: SlideStyle,
    background: Path | None,
) -> Presentation:
    """Render one or more passages into a single presentation."""
    prs = Presentation()
    w, h = style.size_in
    prs.slide_width = Inches(w)
    prs.slide_height = Inches(h)
    for passage in passages:
        # auto-fit the body font per passage so verses fill without overflowing,
        # then guarantee every page fits (shrinking further / erroring if needed)
        tidy_style = fit_body_style(passage.bundles, style)
        passage_style, pages, hang_pt = fit_pages(passage.bundles, tidy_style)
        for page in pages:
            _render_page(prs, page, passage, passage_style, background, hang_pt)
    return prs


def save(prs: Presentation, out_path: str | Path) -> Path:
    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out
